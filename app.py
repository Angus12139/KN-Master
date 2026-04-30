import gradio as gr
from google import genai
from google.genai import types
import os
import requests
import json
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ==========================================
# 1. 基础配置 (采用全新 google.genai 规范)
# ==========================================
client = genai.Client(api_key=os.environ.get("GEMINI_API_KEY"))
FS_APP_ID = os.environ.get("FEISHU_APP_ID")
FS_APP_SECRET = os.environ.get("FEISHU_APP_SECRET")

# ==========================================
# 2. AI 语言处理 (保留原有全部功能)
# ==========================================
prompt_typo = """你的唯一任务是：检查文档或图片中的【中文错别字和语病】。
1. 请保持极高的敏感度，哪怕只有 10% 的错漏把握也请指出，宁可错杀不可放过。
2. 请忽略文档中的纯英文内容，专注中文。
3. 【核心输出格式】：请务必【逐页】输出检查结果。
   - 如果该页没有任何错别字或语病，请严格输出：“第X页：OKOK”
   - 如果该页有需要修改的地方，请清晰列出：“第X页：[原文] -> [修改建议]及原因”"""

prompt_proofread = "做专业的【中英双语校对】。对比中英文翻译是否对齐，检查英文语法和拼写。"
prompt_translate = "做地道的【中译英翻译】。提取内容并输出符合商务规范的纯英文结果。"

def process_ai_task(file_obj, prompt_text):
    if file_obj is None: return "⚠️ 请先上传文件哦！"
    file_name = file_obj.name.lower()
    if not (file_name.endswith('.pdf') or file_name.endswith('.png') or file_name.endswith('.jpg') or file_name.endswith('.jpeg')):
        return "❌ 仅支持 PDF 或图片。请将 PPT/Word 导出为 PDF 后上传。"
    try:
        # 新版 SDK 的文件上传与调用方式
        gemini_file = client.files.upload(file=file_obj.name)
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[prompt_text, gemini_file]
        )
        return response.text
    except Exception as e:
        return f"❌ 错误: {str(e)}"

# ==========================================
# 3. AI 视觉理解核心 (用于飞书图片摘要)
# ==========================================
def get_image_summary(image_bytes):
    """视觉识别 B 列图片并生成 5 字摘要"""
    if not image_bytes: return ""
    try:
        # 新版 SDK 的字节流图片传递方式
        img_part = types.Part.from_bytes(data=image_bytes, mime_type="image/png")
        prompt = "这是演讲幻灯片的下一页内容，请用5个字以内总结其核心要点，作为给演讲者的‘下一页预告’提示（例如：业务增长图表、年度目标展望）"
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[prompt, img_part]
        )
        return response.text.strip()
    except:
        return "图片解析失败"

# ==========================================
# 4. 飞书数据引擎 (读取与回写)
# ==========================================
def get_feishu_token():
    url = "https://open.feishu.cn/open-apis/auth/v3/tenant_access_token/internal"
    r = requests.post(url, json={"app_id": FS_APP_ID, "app_secret": FS_APP_SECRET})
    return r.json().get("tenant_access_token")

def download_fs_media(file_token, token):
    url = f"https://open.feishu.cn/open-apis/drive/v1/medias/{file_token}/download"
    headers = {"Authorization": f"Bearer {token}"}
    r = requests.get(url, headers=headers)
    return r.content if r.status_code == 200 else None

def update_feishu_cell(ss_token, sheet_id, row_index, text, token):
    """【回写引擎】将摘要写回 D 列"""
    col_name = "D"
    range_str = f"{sheet_id}!{col_name}{row_index+1}:{col_name}{row_index+1}"
    url = f"https://open.feishu.cn/open-apis/sheets/v2/spreadsheets/{ss_token}/values"
    headers = {"Authorization": f"Bearer {token}", "Content-Type": "application/json"}
    body = {"valueRange": {"range": range_str, "values": [[text]]}}
    requests.put(url, headers=headers, json=body)

def parse_link(url_link, token):
    headers = {"Authorization": f"Bearer {token}"}
    if "/sheets/" in url_link:
        ss_token = url_link.split("/sheets/")[1].split("?")[0].split("#")[0]
    elif "/wiki/" in url_link:
        wiki_token = url_link.split("/wiki/")[1].split("?")[0].split("#")[0]
        node_res = requests.get(f"https://open.feishu.cn/open-apis/wiki/v2/spaces/get_node?token={wiki_token}", headers=headers).json()
        ss_token = node_res.get("data", {}).get("node", {}).get("obj_token")
    else: return None, None, "格式不支持"
    
    meta_res = requests.get(f"https://open.feishu.cn/open-apis/sheets/v2/spreadsheets/{ss_token}/metainfo", headers=headers).json()
    sheet_id = meta_res["data"]["sheets"][0]["sheetId"]
    return ss_token, sheet_id, "OK"

# ==========================================
# ==========================================
# 5. 引擎 A：自动生成表格摘要并回写 D 列
# ==========================================
def generate_summaries_handler(link):
    if not link: return "⚠️ 请先粘贴链接"
    token = get_feishu_token()
    ss_token, sheet_id, msg = parse_link(link, token)
    if not ss_token: return msg
    
    data_url = f"https://open.feishu.cn/open-apis/sheets/v2/spreadsheets/{ss_token}/values/{sheet_id}!A1:Z500?valueRenderOption=Formula"
    headers = {"Authorization": f"Bearer {token}"}
    raw_data = requests.get(data_url, headers=headers).json().get("data", {}).get("valueRange", {}).get("values", [])
    
    processed_count = 0
    img_col = ord('B') - ord('A')
    
    for i, row in enumerate(raw_data):
        if len(row) > img_col:
            cell_data = row[img_col]
            file_token = None
            
            # 【核心修复】：兼容飞书的各种图片格式（纯图片字典 vs 混合列表）
            if isinstance(cell_data, dict):
                # 如果单元格里只有一张纯图片
                file_token = cell_data.get('fileToken') or cell_data.get('imageToken') or cell_data.get('token')
            elif isinstance(cell_data, list) and len(cell_data) > 0:
                # 如果图片和文字混排，或者被包在列表里
                for item in cell_data:
                    if isinstance(item, dict):
                        file_token = item.get('fileToken') or item.get('imageToken') or item.get('token')
                        if file_token: break # 找到第一个图片就停止
            
            # 如果成功抓到了图片的“身份证号”
            if file_token:
                img_bytes = download_fs_media(file_token, token)
                if img_bytes:
                    summary = get_image_summary(img_bytes)
                    if summary:
                        update_feishu_cell(ss_token, sheet_id, i, summary, token)
                        processed_count += 1
                    
    return f"✅ 摘要生成完毕！已成功在 D 列回写 {processed_count} 条摘要。请刷新飞书表格查看并确认。"
# ==========================================
# 6. 引擎 B：导出智能 PPT (左下角红框，C列正文，D列摘要)
# ==========================================
def export_ppt_handler(link, col_letter):
    if not link: return None, "⚠️ 请先粘贴链接"
    token = get_feishu_token()
    ss_token, sheet_id, msg = parse_link(link, token)
    if not ss_token: return None, msg
    
    headers = {"Authorization": f"Bearer {token}"}
    data_url = f"https://open.feishu.cn/open-apis/sheets/v2/spreadsheets/{ss_token}/values/{sheet_id}!A1:Z500?valueRenderOption=Formula"
    raw_data = requests.get(data_url, headers=headers).json().get("data", {}).get("valueRange", {}).get("values", [])
    
    col_idx = ord(col_letter.upper()) - ord('A')
    hint_col_idx = ord('D') - ord('A')
    
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    
    valid_rows = [r for r in raw_data if len(r) > col_idx]
    
    for i in range(len(valid_rows)):
        row = valid_rows[i]
        
        # A. 正文处理 (C列)
        content_obj = row[col_idx]
        full_text = ""
        segments = []
        if isinstance(content_obj, list):
            segments = content_obj
            for s in segments: full_text += s.get('text', '')
        else:
            full_text = str(content_obj).strip()
            if full_text in ["None", "", "nan"]: continue
            segments = [{'text': full_text, 'segmentStyle': {'foreColor': '#FFFFFF'}}]

        if not full_text.strip(): continue

        # B. 查找下一页提示 (下一行的 D 列)
        next_hint = "演讲结束"
        if i + 1 < len(valid_rows):
            next_row = valid_rows[i+1]
            if len(next_row) > hint_col_idx:
                hint_val = next_row[hint_col_idx]
                if isinstance(hint_val, list):
                    next_hint = "".join([s.get('text','') for s in hint_val]).strip()
                else:
                    next_hint = str(hint_val).strip()
                
                if not next_hint or next_hint == "None":
                    next_hint = "演讲结束"

        # C. 渲染 Slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
        
        margin = Inches(0.8)
        txBox = slide.shapes.add_textbox(margin, margin, prs.slide_width - margin*2, prs.slide_height - Inches(2))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = 1 
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        
        text_len = len(full_text)
        font_size = 80 if text_len <= 20 else 60 if text_len <= 60 else 40
        
        for seg in segments:
            run = p.add_run()
            run.text = seg.get('text', '')
            run.font.name, run.font.size = '微软雅黑', Pt(font_size)
            run.font.bold = seg.get('segmentStyle', {}).get('bold', True)
            c = seg.get('segmentStyle', {}).get('foreColor', '#FFFFFF')
            if c.upper() in ["#000000", "#121212"]: run.font.color.rgb = RGBColor(255, 255, 255)
            else:
                try: 
                    hex_c = c.lstrip('#')
                    run.font.color.rgb = RGBColor(*(int(hex_c[k:k+2], 16) for k in (0, 2, 4)))
                except: run.font.color.rgb = RGBColor(255, 255, 255)

        # D. 左下角红色预告框
        box_w, box_h = Inches(5.0), Inches(0.6)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, prs.slide_height - Inches(1.2), box_w, box_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(225, 29, 72)
        shape.line.fill.background()
        
        tf_h = shape.text_frame
        tf_h.vertical_anchor = 1
        ph = tf_h.paragraphs[0]
        ph.alignment = PP_ALIGN.LEFT
        rh = ph.add_run()
        rh.text = f" 下一页预告：{next_hint}"
        rh.font.name, rh.font.size, rh.font.bold = '微软雅黑', Pt(20), True
        rh.font.color.rgb = RGBColor(255, 255, 255)

    out = "V5.6_智能预测提词器.pptx"
    prs.save(out)
    return out, f"✅ 已成功导出 {len(valid_rows)} 页带下一页预告的 PPT！"

# ==========================================
# 7. UI 界面整合
# ==========================================
with gr.Blocks(title="AI 智能文档工作站 V5.6") as demo:
    gr.Markdown("# 🚀 AI 智能文档工作站 V5.6 (全能版)")
    
    with gr.Tabs():
        # --- 保留的 Tab 1：纠错校对 ---
        with gr.TabItem("📝 AI 语言处理"):
            with gr.Row():
                with gr.Column():
                    ai_file = gr.File(label="上传 PDF 或图片")
                    with gr.Row():
                        btn_typo = gr.Button("🇨🇳 逐页纠错", variant="primary")
                        btn_proof = gr.Button("⚖️ 双语校对", variant="secondary")
                        btn_trans = gr.Button("🌍 翻译结果", variant="secondary")
                ai_output = gr.Textbox(label="AI 处理结果", lines=20)
            
            btn_typo.click(fn=lambda f: process_ai_task(f, prompt_typo), inputs=ai_file, outputs=ai_output)
            btn_proof.click(fn=lambda f: process_ai_task(f, prompt_proofread), inputs=ai_file, outputs=ai_output)
            btn_trans.click(fn=lambda f: process_ai_task(f, prompt_translate), inputs=ai_file, outputs=ai_output)

        # --- 全新的 Tab 2：飞书双引擎 ---
        with gr.TabItem("🎬 飞书一键转 PPT"):
            gr.Markdown("### 双步工作流：1. 读取 B 列图片生成摘要写回 D 列 ➡️ 2. 读取 C 列正文与下页 D 列生成 PPT")
            with gr.Row():
                link_input = gr.Textbox(label="第一步：粘贴飞书链接 (Sheet/Wiki)", placeholder="https://...")
                col_input = gr.Textbox(label="正文所在列 (默认 C 列)", value="C")

            with gr.Row():
                with gr.Column():
                    gr.Markdown("### 🛠 引擎 A：表格内容生产")
                    summary_btn = gr.Button("🤖 识别图片生成提示词 (将结果写入 D 列)", variant="secondary")
                    summary_status = gr.Markdown("状态：等待指令")
                
                with gr.Column():
                    gr.Markdown("### 🚀 引擎 B：PPT 智能导出")
                    export_btn = gr.Button("🔥 立即导出智能预测 PPT", variant="primary")
                    ppt_file = gr.File(label="下载导出的 PPT")
                    export_status = gr.Markdown("状态：等待指令")

            summary_btn.click(fn=generate_summaries_handler, inputs=link_input, outputs=summary_status)
            export_btn.click(fn=export_ppt_handler, inputs=[link_input, col_input], outputs=[ppt_file, export_status])

demo.launch()